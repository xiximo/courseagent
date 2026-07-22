"""评审报告 PPT 导出。"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas.review_report import ReviewReportDto

MAX_TABLE_ROWS = 6


def render_review_report_pptx(report: ReviewReportDto) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    _add_title_slide(presentation, report)
    _add_basic_info_slide(presentation, report)
    _add_missing_features_slide(presentation, report)
    _add_improvements_slide(presentation, report)
    _add_summary_slide(presentation, report)
    _add_conclusion_slide(presentation, report)

    if report.citations:
        _add_citations_slide(presentation, report)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _add_title_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = report.title
    subtitle = slide.placeholders[1]
    info = report.basicInfo
    subtitle.text = "\n".join(
        part
        for part in [
            info.standardName,
            f"标准编号：{info.standardNo or '—'}",
            f"评审日期：{info.reviewDate or '—'}",
        ]
        if part
    )


def _add_basic_info_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = _add_content_slide(presentation, "1. 基本信息")
    info = report.basicInfo
    lines = [
        f"标准名称：{info.standardName or '—'}",
        f"标准编号：{info.standardNo or '—'}",
        f"所属零部件：{info.partName or '—'}",
        f"评审日期：{info.reviewDate or '—'}",
    ]
    _fill_text_box(slide, lines)


def _add_missing_features_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = _add_content_slide(presentation, "2. 缺失特性结论")
    if not report.missingFeatures:
        _fill_text_box(slide, ["未识别到缺失特性"])
        return

    rows = [["缺失特性说明", "涉及标准"]]
    rows.extend(
        [row.description, row.relatedStandards or "—"]
        for row in report.missingFeatures[:MAX_TABLE_ROWS]
    )
    _add_table(slide, rows)
    if len(report.missingFeatures) > MAX_TABLE_ROWS:
        _append_note(slide, "更多条目请查看 Word 完整报告")


def _add_improvements_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = _add_content_slide(presentation, "3. 需提升结论")
    if not report.improvements:
        _fill_text_box(slide, ["未识别到需提升项"])
        return

    rows = [["特性要素", "建议内容", "建议方法", "参考标准"]]
    rows.extend(
        [
            row.featureElement,
            row.suggestedContent or "—",
            row.suggestedMethod or "—",
            row.referenceStandards or "—",
        ]
        for row in report.improvements[:MAX_TABLE_ROWS]
    )
    _add_table(slide, rows, col_widths=[1.4, 2.2, 2.2, 2.0])
    if len(report.improvements) > MAX_TABLE_ROWS:
        _append_note(slide, "更多条目请查看 Word 完整报告")


def _add_summary_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = _add_content_slide(presentation, "4. 评审总结建议")
    if not report.summaryPoints:
        _fill_text_box(slide, ["暂无总结建议"])
        return

    lines: list[str] = []
    for point in report.summaryPoints[:8]:
        if point.title and point.content:
            lines.append(f"{point.order}. {point.title}：{point.content}")
        elif point.title:
            lines.append(f"{point.order}. {point.title}")
        else:
            lines.append(f"{point.order}. {point.content}")
    _fill_text_box(slide, lines)


def _add_conclusion_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = _add_content_slide(presentation, "5. 最终结论")
    _fill_text_box(slide, [report.finalConclusion or "（无）", "", report.disclaimer])


def _add_citations_slide(presentation: Presentation, report: ReviewReportDto) -> None:
    slide = _add_content_slide(presentation, "引用依据")
    lines = []
    for cite in report.citations[:8]:
        line = " ".join(
            part
            for part in [cite.standardNo, cite.standardName, cite.position]
            if part
        )
        if line:
            lines.append(f"• {line}")
    _fill_text_box(slide, lines or ["（无）"])


def _add_content_slide(presentation: Presentation, title: str):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    return slide


def _fill_text_box(slide, lines: list[str]) -> None:
    body = slide.shapes.placeholders[1].text_frame
    body.text = lines[0] if lines else ""
    for line in lines[1:]:
        paragraph = body.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(18)


def _add_table(
    slide,
    rows: list[list[str]],
    *,
    col_widths: list[float] | None = None,
) -> None:
    if not rows:
        return

    row_count = len(rows)
    col_count = len(rows[0])
    left = Inches(0.4)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(0.4 + 0.35 * row_count)

    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for index, inch_width in enumerate(col_widths):
            table.columns[index].width = Inches(inch_width)

    for row_index, row_values in enumerate(rows):
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12 if row_index == 0 else 11)
                if row_index == 0:
                    paragraph.font.bold = True


def _append_note(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.0), Inches(0.5))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(12)
    paragraph.font.italic = True
    paragraph.alignment = PP_ALIGN.RIGHT
